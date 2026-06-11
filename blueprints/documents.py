import os
import uuid
from flask import Blueprint, render_template, request, redirect, url_for, flash, send_file, jsonify
from flask_login import login_required, current_user
from models import db, Document, DocumentCategory
from config import Config
from datetime import datetime
from blueprints.auth import admin_required
from werkzeug.utils import secure_filename

documents_bp = Blueprint('documents', __name__)


def allowed_file(filename):
    basename = os.path.basename(filename)
    return '.' in basename and basename.rsplit('.', 1)[1].lower() in Config.ALLOWED_EXTENSIONS


def safe_filename(filename):
    basename = os.path.basename(filename)
    if not basename or basename == '':
        basename = filename.replace('/', '_').replace('\\', '_')
    return basename


def get_file_type(filename):
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    image_types = {'jpg', 'jpeg', 'png', 'gif', 'bmp', 'svg'}
    doc_types = {'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx'}
    cad_types = {'dwg', 'dxf', 'step', 'stp', 'igs', 'iges'}
    archive_types = {'zip', 'rar', '7z', 'tar', 'gz'}
    if ext in image_types:
        return 'image'
    elif ext in doc_types:
        return 'document'
    elif ext in cad_types:
        return 'cad'
    elif ext in archive_types:
        return 'archive'
    return 'other'


def search_filesystem(search_term, base_path):
    """在文件系统中搜索匹配的文件（优化版）"""
    results = []
    search_lower = search_term.lower()
    
    # 需要跳过的目录（包含大量CAD/图纸文件，搜索时跳过避免超时）
    skip_dirs = {'MCS图纸_DXF测试', 'MCS图纸_DXF_PDF', 'BOM_清单', '__pycache__', '.git'}
    
    for root, dirs, files in os.walk(base_path):
        # 过滤掉大型CAD目录
        dirs[:] = [d for d in dirs if not d.startswith('.') and d not in skip_dirs]
        
        for filename in files:
            if filename.startswith('.'):
                continue
            
            # 检查文件名是否匹配
            if search_lower in filename.lower():
                filepath = os.path.join(root, filename)
                
                try:
                    stat = os.stat(filepath)
                    size = stat.st_size
                    if size < 1024:
                        size_str = f'{size}B'
                    elif size < 1024*1024:
                        size_str = f'{size/1024:.1f}KB'
                    elif size < 1024*1024*1024:
                        size_str = f'{size/(1024*1024):.1f}MB'
                    else:
                        size_str = f'{size/(1024*1024*1024):.1f}GB'
                    
                    rel_path = os.path.relpath(filepath, base_path)
                    results.append({
                        'type': 'filesystem',
                        'name': filename,
                        'path': rel_path.replace('\\', '/'),
                        'size': size_str,
                        'ext': filename.rsplit('.', 1)[1].lower() if '.' in filename else '',
                        'static_url': url_for('static', filename='uploads/' + rel_path.replace('\\', '/')),
                    })
                    
                    # 限制搜索结果数量
                    if len(results) >= 100:
                        return results
                except Exception:
                    continue
    
    results.sort(key=lambda x: x['name'].lower())
    return results


@documents_bp.route('/')
@login_required
def index():
    category_id = request.args.get('category_id', type=int)
    search = request.args.get('search', '').strip()
    categories = DocumentCategory.query.order_by(DocumentCategory.sort_order).all()

    query = Document.query.order_by(Document.created_at.desc())
    if category_id:
        query = query.filter_by(category_id=category_id)
    
    # 数据库文档搜索
    db_documents = []
    if search:
        query = query.filter(Document.title.contains(search) | Document.original_filename.contains(search))
        db_documents = query.all()
    else:
        db_documents = query.all()
    
    # 文件系统搜索（仅在有搜索关键字时）
    fs_documents = []
    if search:
        fs_documents = search_filesystem(search, Config.UPLOAD_FOLDER)
    
    return render_template('documents/index.html',
                           documents=db_documents,
                           fs_documents=fs_documents,
                           categories=categories,
                           current_category_id=category_id,
                           search=search)


@documents_bp.route('/upload', methods=['POST'])
@login_required
def upload():
    category_id = request.form.get('category_id', type=int)
    title = request.form.get('title', '').strip()
    description = request.form.get('description', '').strip()
    files = request.files.getlist('files')

    if not category_id or not files or not files[0].filename:
        flash('请选择分类并上传文件', 'error')
        return redirect(url_for('documents.index'))

    category = DocumentCategory.query.get(category_id)
    if not category:
        flash('分类不存在', 'error')
        return redirect(url_for('documents.index'))

    count = 0
    for file in files:
        if file and allowed_file(file.filename):
            orig_name = safe_filename(file.filename)
            unique_name = f"{uuid.uuid4().hex}_{orig_name}"
            filepath = os.path.join(Config.UPLOAD_FOLDER, unique_name)
            file.save(filepath)
            file_size = os.path.getsize(filepath)

            t = title if title else os.path.splitext(orig_name)[0]
            doc = Document(
                category_id=category_id,
                title=t,
                filename=unique_name,
                original_filename=orig_name,
                file_size=file_size,
                file_type=get_file_type(file.filename),
                uploader_id=current_user.id,
                description=description
            )
            db.session.add(doc)
            count += 1

    db.session.commit()
    flash(f'成功上传 {count} 个文件', 'success')
    return redirect(url_for('documents.index', category_id=category_id))


@documents_bp.route('/download/<int:doc_id>')
@login_required
def download(doc_id):
    doc = Document.query.get_or_404(doc_id)
    doc.download_count += 1
    db.session.commit()
    filepath = os.path.join(Config.UPLOAD_FOLDER, doc.filename)
    if not os.path.exists(filepath):
        flash('文件不存在', 'error')
        return redirect(url_for('documents.index'))
    return send_file(filepath, download_name=doc.original_filename, as_attachment=True)


@documents_bp.route('/view/<int:doc_id>')
@login_required
def view(doc_id):
    doc = Document.query.get_or_404(doc_id)
    filepath = os.path.join(Config.UPLOAD_FOLDER, doc.filename)
    if not os.path.exists(filepath):
        flash('文件不存在', 'error')
        return redirect(url_for('documents.index'))
    image_types = {'jpg', 'jpeg', 'png', 'gif', 'bmp', 'svg'}
    ext = doc.original_filename.rsplit('.', 1)[1].lower() if '.' in doc.original_filename else ''
    if ext in image_types:
        return send_file(filepath)
    return send_file(filepath, download_name=doc.original_filename)


@documents_bp.route('/delete/<int:doc_id>', methods=['POST'])
@login_required
def delete(doc_id):
    doc = Document.query.get_or_404(doc_id)
    if not current_user.is_admin and doc.uploader_id != current_user.id:
        flash('无权删除', 'error')
        return redirect(url_for('documents.index'))
    filepath = os.path.join(Config.UPLOAD_FOLDER, doc.filename)
    if os.path.exists(filepath):
        os.remove(filepath)
    db.session.delete(doc)
    db.session.commit()
    flash('文件已删除', 'success')
    return redirect(url_for('documents.index'))


@documents_bp.route('/categories', methods=['GET', 'POST'])
@login_required
@admin_required
def manage_categories():
    if request.method == 'POST':
        name = request.form.get('name', '').strip()
        description = request.form.get('description', '').strip()
        sort_order = request.form.get('sort_order', 0, type=int)
        if name:
            cat = DocumentCategory(name=name, description=description, sort_order=sort_order)
            db.session.add(cat)
            db.session.commit()
            flash('分类添加成功', 'success')
        return redirect(url_for('documents.manage_categories'))
    categories = DocumentCategory.query.order_by(DocumentCategory.sort_order).all()
    return render_template('documents/categories.html', categories=categories)


@documents_bp.route('/categories/delete/<int:cat_id>', methods=['POST'])
@login_required
@admin_required
def delete_category(cat_id):
    cat = DocumentCategory.query.get_or_404(cat_id)
    if cat.documents.count() > 0:
        flash('该分类下还有文件，无法删除', 'error')
        return redirect(url_for('documents.manage_categories'))
    db.session.delete(cat)
    db.session.commit()
    flash('分类已删除', 'success')
    return redirect(url_for('documents.manage_categories'))


@documents_bp.route('/browse/')
@documents_bp.route('/browse/<path:subpath>')
@login_required
def browse(subpath=''):
    base = Config.UPLOAD_FOLDER
    safe_sub = subpath.replace('\\', '/').strip('/')
    current_dir = os.path.join(base, safe_sub) if safe_sub else base

    real_base = os.path.realpath(base)
    real_current = os.path.realpath(current_dir)
    if not real_current.startswith(real_base + os.sep) and real_current != real_base:
        flash('路径不允许', 'error')
        return redirect(url_for('documents.browse'))

    folders = []
    files = []
    try:
        for entry in sorted(os.scandir(current_dir), key=lambda e: (not e.is_dir(), e.name.lower())):
            if entry.name.startswith('.'):
                continue
            rel = (safe_sub + '/' + entry.name).strip('/') if safe_sub else entry.name
            if entry.is_dir():
                folders.append({'name': entry.name, 'path': rel})
            elif entry.is_file():
                stat = entry.stat()
                size = stat.st_size
                if size < 1024:
                    size_str = f'{size}B'
                elif size < 1024*1024:
                    size_str = f'{size/1024:.1f}KB'
                elif size < 1024*1024*1024:
                    size_str = f'{size/(1024*1024):.1f}MB'
                else:
                    size_str = f'{size/(1024*1024*1024):.1f}GB'
                files.append({
                    'name': entry.name,
                    'path': rel,
                    'size': size_str,
                    'ext': entry.name.rsplit('.', 1)[1].lower() if '.' in entry.name else '',
                    'static_url': url_for('static', filename='uploads/' + rel.replace('\\', '/')),
                })
    except PermissionError:
        flash('无权限访问', 'error')
        return redirect(url_for('documents.browse'))

    breadcrumbs = []
    if safe_sub:
        parts = safe_sub.split('/')
        for i in range(len(parts)):
            breadcrumbs.append({
                'name': parts[i],
                'path': '/'.join(parts[:i+1]),
            })

    return render_template('documents/browse.html',
                           folders=folders, files=files,
                           subpath=safe_sub, breadcrumbs=breadcrumbs)


@documents_bp.route('/raw/<int:doc_id>')
@login_required
def raw_file(doc_id):
    """返回DB文档的原始文件内容，供内置查看器使用"""
    doc = Document.query.get_or_404(doc_id)
    filepath = os.path.join(Config.UPLOAD_FOLDER, doc.filename)
    if not os.path.exists(filepath):
        flash('文件不存在', 'error')
        return redirect(url_for('documents.index'))
    return send_file(filepath, download_name=doc.original_filename)


@documents_bp.route('/pptx_preview/<int:doc_id>')
@login_required
def pptx_preview(doc_id):
    """提取PPTX幻灯片文本为HTML"""
    try:
        from pptx import Presentation
    except ImportError:
        return '<p style="color:red;">python-pptx未安装</p>', 500

    doc = Document.query.get_or_404(doc_id)
    filepath = os.path.join(Config.UPLOAD_FOLDER, doc.filename)
    if not os.path.exists(filepath):
        return '<p>文件不存在</p>', 404

    try:
        prs = Presentation(filepath)
        html = '<div style="font-family:sans-serif;">'
        for i, slide in enumerate(prs.slides, 1):
            html += f'<div style="margin:10px 0;padding:16px;background:#fff;border-radius:12px;box-shadow:0 2px 8px rgba(0,0,0,0.1);">'
            html += f'<div style="font-size:14px;font-weight:700;color:#00897b;margin-bottom:8px;">📊 幻灯片 {i}/{len(prs.slides)}</div>'
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            font_size = '13px'
                            is_bold = False
                            for run in para.runs:
                                if run.font.size:
                                    font_size = f'{run.font.size.pt}pt'
                                if run.font.bold:
                                    is_bold = True
                            style = f'font-size:{font_size};'
                            if is_bold:
                                style += 'font-weight:700;'
                            html += f'<p style="{style}margin:4px 0;">{text}</p>'
                if shape.has_table:
                    table = shape.table
                    html += '<table style="border-collapse:collapse;width:100%;margin:8px 0;font-size:12px;">'
                    for row in table.rows:
                        html += '<tr>'
                        for cell in row.cells:
                            html += f'<td style="border:1px solid #ddd;padding:4px 8px;">{cell.text}</td>'
                        html += '</tr>'
                    html += '</table>'
            html += '</div>'
        html += '</div>'
        return html
    except Exception as e:
        return f'<p style="color:red;">PPTX解析失败: {str(e)}</p>', 500


@documents_bp.route('/viewer')
@login_required
def viewer():
    """内置文件查看器 - 支持PDF/Excel/DOCX/图片/文本预览"""
    doc_id = request.args.get('doc_id', '')
    file_url = request.args.get('file', '')
    title = request.args.get('title', '')

    # 确定文件URL和类型
    if doc_id:
        doc = Document.query.get(int(doc_id))
        if not doc:
            return '文档不存在', 404
        file_url = url_for('documents.raw_file', doc_id=doc_id)
        title = title or doc.original_filename
        ext = os.path.splitext(doc.original_filename)[1].lower()
    elif file_url:
        title = title or file_url.split('/')[-1]
        ext = os.path.splitext(title)[1].lower() if '.' in title else ''
        # 文件系统文件转换为static URL
        if not file_url.startswith('/'):
            file_url = url_for('static', filename='uploads/' + file_url)
    else:
        return '缺少文件参数', 400

    # 判断文件类型
    image_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp'}
    pdf_ext = '.pdf'
    excel_exts = {'.xls', '.xlsx', '.xlsm', '.csv'}
    doc_exts = {'.doc', '.docx'}
    text_exts = {'.txt', '.log', '.md', '.py', '.html', '.css', '.js', '.json', '.xml', '.scr'}

    viewer_type = 'unknown'
    pptx_preview_html = ''
    
    if ext in image_exts:
        viewer_type = 'image'
    elif ext == pdf_ext:
        viewer_type = 'pdf'
    elif ext in excel_exts:
        viewer_type = 'excel'
    elif ext == '.docx':
        viewer_type = 'doc'
    elif ext == '.doc':
        viewer_type = 'olddoc'
    elif ext == '.pptx':
        # 服务端提取PPTX文本
        try:
            from pptx import Presentation
            if doc_id:
                doc = Document.query.get(int(doc_id))
                fpath = os.path.join(Config.UPLOAD_FOLDER, doc.filename)
            else:
                fpath = os.path.join(Config.UPLOAD_FOLDER, file_url.replace('\\', '/'))
            if os.path.exists(fpath):
                prs = Presentation(fpath)
                h = '<div style="font-family:sans-serif;">'
                for i, slide in enumerate(prs.slides, 1):
                    h += f'<div style="margin:10px 0;padding:16px;background:#fff;border-radius:12px;box-shadow:0 2px 8px rgba(0,0,0,0.1);">'
                    h += f'<div style="font-size:14px;font-weight:700;color:#00897b;margin-bottom:8px;">📊 幻灯片 {i}/{len(prs.slides)}</div>'
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    fs = '13px'
                                    bd = ''
                                    for run in para.runs:
                                        if run.font.size: fs = f'{run.font.size.pt}pt'
                                        if run.font.bold: bd = 'font-weight:700;'
                                    h += f'<p style="font-size:{fs};{bd}margin:4px 0;">{text}</p>'
                        if shape.has_table:
                            h += '<table style="border-collapse:collapse;width:100%;margin:8px 0;font-size:12px;">'
                            for row in shape.table.rows:
                                h += '<tr>'
                                for cell in row.cells:
                                    h += f'<td style="border:1px solid #ddd;padding:4px 8px;">{cell.text}</td>'
                                h += '</tr>'
                            h += '</table>'
                    h += '</div>'
                h += '</div>'
                pptx_preview_html = h
                viewer_type = 'pptx_done'
        except Exception as e:
            pptx_preview_html = f'<div class="unsupported"><p>PPTX解析失败: {str(e)}</p><a href="{file_url}" download style="display:inline-block;margin-top:12px;padding:8px 24px;background:#00897b;color:#fff;border-radius:20px;text-decoration:none;">下载文件</a></div>'
            viewer_type = 'pptx_done'
    elif ext == '.ppt':
        viewer_type = 'oldppt'
    elif ext in text_exts:
        viewer_type = 'text'

    return f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>查看: {title}</title>
     <style>
         * {{ margin: 0; padding: 0; box-sizing: border-box; }}
         body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #f0f0f0; }}
         .viewer-header {{
             display: flex; align-items: center; gap: 12px; padding: 8px 12px;
             background: #1a1a2e; color: #fff; position: sticky; top: 0; z-index: 100;
         }}
         .viewer-header .back {{ color: #fff; text-decoration: none; font-size: 20px; }}
         .viewer-header .title {{ flex: 1; font-size: 13px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }}
         .viewer-header .actions {{ display: flex; gap: 8px; }}
         .viewer-header .actions a {{
             padding: 4px 10px; border-radius: 12px; font-size: 11px; text-decoration: none;
             background: rgba(255,255,255,0.2); color: #fff;
         }}
         .viewer-content {{ width: 100%; }}
         #viewer {{ width: 100%; overflow: auto; }}
         .page-container {{ margin: 0 auto 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.2); position: relative; }}
         canvas {{ display: block; margin: 0 auto; }}
         .image-viewer {{ text-align: center; padding: 10px; }}
         .image-viewer img {{ max-width: 100%; height: auto; }}
         .text-viewer {{ padding: 12px; background: #fff; font-family: monospace; font-size: 13px; line-height: 1.5; white-space: pre-wrap; word-break: break-all; min-height: 100vh; }}
         .excel-viewer {{ padding: 10px; background: #fff; overflow: auto; }}
         .excel-viewer table {{ border-collapse: collapse; font-size: 12px; }}
         .excel-viewer td, .excel-viewer th {{ border: 1px solid #ddd; padding: 6px 8px; min-width: 60px; }}
         .excel-viewer th {{ background: #f5f5f5; font-weight: 600; position: sticky; top: 0; }}
         .doc-viewer {{ padding: 16px; background: #fff; max-width: 800px; margin: 0 auto; min-height: 100vh; font-size: 14px; line-height: 1.8; }}
         .loading {{ text-align: center; padding: 60px; color: #999; }}
         .unsupported {{ text-align: center; padding: 60px; color: #999; }}
         .unsupported i {{ font-size: 48px; display: block; margin-bottom: 10px; color: #ddd; }}
     </style>
 </head>
 <body>
 <div class="viewer-header">
     <a href="javascript:history.back()" class="back">←</a>
     <span class="title">{title}</span>
     <div class="actions">
         <a href="{file_url}" download>下载</a>
     </div>
 </div>
 <div class="viewer-content">
     <div id="viewer">
         <div id="loadingArea" class="loading">加载中...</div>
     </div>
 </div>
 <script>
 const FILE_URL = "{file_url}";
 const VIEWER_TYPE = "{viewer_type}";
 const TITLE = "{title}";
 const FILE_EXT = "{ext}";

function hideLoading() {{ const el = document.getElementById('loadingArea'); if (el) el.style.display = 'none'; }}

async function initViewer() {{
    const viewer = document.getElementById('viewer');

    if (VIEWER_TYPE === 'image') {{
        viewer.innerHTML = `<div class="image-viewer"><img src="${{FILE_URL}}" alt="${{TITLE}}"></div>`;
        hideLoading();
    }}
    else if (VIEWER_TYPE === 'text') {{
        try {{
            const resp = await fetch(FILE_URL);
            const text = await resp.text();
            viewer.innerHTML = `<div class="text-viewer">${{escapeHtml(text)}}</div>`;
        }} catch(e) {{
            viewer.innerHTML = '<div class="loading">加载失败: ' + e.message + '</div>';
        }}
        hideLoading();
    }}
    else if (VIEWER_TYPE === 'pdf') {{
        // PDF.js CDN
        const script = document.createElement('script');
        script.src = 'https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js';
        script.onload = () => renderPDF();
        document.head.appendChild(script);
        const workerScript = document.createElement('script');
        workerScript.src = 'https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js';
        document.head.appendChild(workerScript);
    }}
    else if (VIEWER_TYPE === 'excel') {{
        const script = document.createElement('script');
        script.src = 'https://cdnjs.cloudflare.com/ajax/libs/xlsx/0.18.5/xlsx.full.min.js';
        script.onload = () => renderExcel();
        document.head.appendChild(script);
    }}
    else if (VIEWER_TYPE === 'doc') {{
         const script = document.createElement('script');
         script.src = 'https://cdnjs.cloudflare.com/ajax/libs/mammoth/1.6.0/mammoth.browser.min.js';
         script.onload = () => renderDoc();
         document.head.appendChild(script);
     }}
     else if (VIEWER_TYPE === 'pptx_done') {{
         // 服务端已渲染
         viewer.innerHTML = `{pptx_preview_html}`;
         hideLoading();
     }}
     else if (VIEWER_TYPE === 'olddoc' || VIEWER_TYPE === 'oldppt') {{
         const typeName = VIEWER_TYPE === 'olddoc' ? '旧版 .doc' : '旧版 .ppt';
         viewer.innerHTML = `<div class="unsupported"><i>📄</i><p>${{typeName}} 格式不支持在线预览</p><p style="font-size:12px;color:#ccc;">请下载后用 Office/WPS 打开</p><a href="${{FILE_URL}}" download style="display:inline-block;margin-top:12px;padding:8px 24px;background:#00897b;color:#fff;border-radius:20px;text-decoration:none;">下载文件</a></div>`;
         hideLoading();
     }}
     else {{
          viewer.innerHTML = `<div class="unsupported"><i>📄</i><p>不支持预览此文件类型</p><p style="font-size:12px;color:#ccc;">${{TITLE}}</p><a href="${{FILE_URL}}" download style="display:inline-block;margin-top:12px;padding:8px 24px;background:#00897b;color:#fff;border-radius:20px;text-decoration:none;">下载文件</a></div>`;
          hideLoading();
      }}
}}

async function renderPDF() {{
    try {{
        const pdf = await pdfjsLib.getDocument(FILE_URL).promise;
        const viewer = document.getElementById('viewer');
        viewer.innerHTML = '';
        const dpr = window.devicePixelRatio || 1;
        for (let pageNum = 1; pageNum <= pdf.numPages; pageNum++) {{
            const page = await pdf.getPage(pageNum);
            const fitScale = (window.innerWidth - 16) / page.getViewport({{ scale: 1 }}).width;
            const hiresScale = Math.max(fitScale, 1.2) * dpr;
            const fitVp = page.getViewport({{ scale: fitScale * dpr }});
            const hiresVp = page.getViewport({{ scale: hiresScale }});
            
            const canvas = document.createElement('canvas');
            const ctx = canvas.getContext('2d');
            canvas.width = fitVp.width;
            canvas.height = fitVp.height;
            canvas.style.width = (fitVp.width / dpr) + 'px';
            canvas.style.height = (fitVp.height / dpr) + 'px';
            canvas.style.cursor = 'zoom-in';
            canvas.title = '点击放大查看';
            
            const container = document.createElement('div');
            container.className = 'page-container';
            container.appendChild(canvas);
            viewer.appendChild(container);
            
            // 首次渲染：适应屏幕
            await page.render({{ canvasContext: ctx, viewport: fitVp }}).promise;
            
            let zoomed = false;
            canvas.onclick = async function() {{
                if (!zoomed) {{
                    canvas.width = hiresVp.width;
                    canvas.height = hiresVp.height;
                    canvas.style.width = (hiresVp.width / dpr) + 'px';
                    canvas.style.height = (hiresVp.height / dpr) + 'px';
                    canvas.style.cursor = 'zoom-out';
                    canvas.title = '点击缩小';
                    await page.render({{ canvasContext: ctx, viewport: hiresVp }}).promise;
                    zoomed = true;
                }} else {{
                    canvas.width = fitVp.width;
                    canvas.height = fitVp.height;
                    canvas.style.width = (fitVp.width / dpr) + 'px';
                    canvas.style.height = (fitVp.height / dpr) + 'px';
                    canvas.style.cursor = 'zoom-in';
                    canvas.title = '点击放大查看';
                    await page.render({{ canvasContext: ctx, viewport: fitVp }}).promise;
                    zoomed = false;
                }}
            }};
        }}
        hideLoading();
    }} catch(e) {{
        document.getElementById('viewer').innerHTML = '<div class="loading">PDF加载失败: ' + e.message + '</div>';
        hideLoading();
    }}
}}

async function renderExcel() {{
    try {{
        const resp = await fetch(FILE_URL);
        const data = await resp.arrayBuffer();
        let wb;
        if (FILE_EXT === '.xls') {{
            // 旧版.xls用binary方式读取
            const bytes = new Uint8Array(data);
            let binary = '';
            for (let i = 0; i < bytes.byteLength; i++) binary += String.fromCharCode(bytes[i]);
            wb = XLSX.read(binary, {{ type: 'binary' }});
        }} else {{
            wb = XLSX.read(data, {{ type: 'array' }});
        }}
        let html = '<div class="excel-viewer">';
        wb.SheetNames.forEach((name) => {{
            const ws = wb.Sheets[name];
            const tableHtml = XLSX.utils.sheet_to_html(ws, {{ editable: false }});
            if (wb.SheetNames.length > 1) {{
                html += `<h3 style="margin:8px 0;font-size:13px;">${{name}}</h3>`;
            }}
            html += tableHtml;
        }});
        html += '</div>';
        document.getElementById('viewer').innerHTML = html;
        hideLoading();
    }} catch(e) {{
        document.getElementById('viewer').innerHTML = '<div class="loading">Excel加载失败: ' + e.message + '</div>';
        hideLoading();
    }}
}}

async function renderDoc() {{
    if (FILE_EXT === '.doc') {{
        document.getElementById('viewer').innerHTML = `<div class="unsupported"><i>📄</i><p>旧版 .doc 格式不支持在线预览</p><p style="font-size:12px;color:#ccc;">请下载后用 Word/WPS 打开</p><a href="${{FILE_URL}}" download style="display:inline-block;margin-top:12px;padding:8px 24px;background:#00897b;color:#fff;border-radius:20px;text-decoration:none;">下载文件</a></div>`;
        hideLoading();
        return;
    }}
    try {{
        const resp = await fetch(FILE_URL);
        const data = await resp.arrayBuffer();
        const result = await mammoth.convertToHtml({{ arrayBuffer: data }});
        document.getElementById('viewer').innerHTML = `<div class="doc-viewer">${{result.value}}</div>`;
        hideLoading();
    }} catch(e) {{
        document.getElementById('viewer').innerHTML = '<div class="loading">文档加载失败: ' + e.message + '</div>';
        hideLoading();
    }}
}}

function escapeHtml(text) {{
    const d = document.createElement('div');
    d.textContent = text;
    return d.innerHTML;
}}

initViewer();
</script>
</body>
</html>'''
